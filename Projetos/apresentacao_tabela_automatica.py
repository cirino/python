from pptx import Presentation
from pptx.util import Inches
import pyexcel as pe

print("""
    Exemplo de criação de apresentação PPTX em loop utilizando dados de Excel
    Vish, o bagulho foi loko pra conseguir criar este aplicativo mano
        -> agora aprendi, já era
    Day 24 Code Python - 23/05/2018
""")

dadosExcel = pe.iget_records(file_name="apresentacao_automatica.xlsx")  # tentar criar uma função

prs = Presentation()    # se for ler um PPTX, passar como parâmetro
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = 'Idades -> Feito com Python'

# configuração de colunas e celulas --> ainda preciso descobrir como automatizar estes tamanhos porque o bagulho tá loko
rows = 5
cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

# tamanho das colunas
table.columns[0].width = Inches(3.0)
table.columns[1].width = Inches(2.0)

# nome das colunas - fixo na posição 0
table.cell(0, 0).text = 'Nome'
table.cell(0, 1).text = 'Idade'

nome = []
idade = []
cont = 1

for itens in dadosExcel:
    # escrevendo os dados na células
    if cont > 0:    # o 0 sempre será o header, e este é fixo.
        table.cell(cont, 0).text = str(itens['nome'])
        table.cell(cont, 1).text = str(itens['idade'])

    cont += 1

# liberando o recurso, iiirraaaaa
pe.free_resources()

# salvando o arquivo pptx
prs.save('apresentacao_tabela_automatica.pptx')

print('-' * 34)
print('APRESENTAÇÃO CRIADA COM SUCESSO.')
