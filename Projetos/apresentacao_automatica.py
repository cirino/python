import pyexcel as pe
from pptx import Presentation

print("""
    Exemplo de criação de apresentação PPTX simples utilizando dados de Excel
    Day 24 Code Python - 23/05/2018
""")

records = pe.iget_records(file_name="/mnt/VersaoBeta/Python/PythonBasico/Projetos/apresentacao_automatica.xlsx")

texto = ['python-pptx funcionou legal!']

for record in records:
    texto.append(str(record['nome']) + ' com idade ' + str(record['idade']))

pe.free_resources()

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Yes mano!"
subtitle.text = str(texto)

prs.save('/mnt/VersaoBeta/Python/PythonBasico/Projetos/apresentacao_automatica.pptx')

print(texto)
